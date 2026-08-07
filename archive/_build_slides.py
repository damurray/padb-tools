"""One-off generator for PADB_Simple_and_Interactive.pptx. Not part of the
padb-tools pipeline -- run once, then archive alongside the other one-off
scripts if no longer needed."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IMG_DIR = r"C:\apps\padb\tools\_tutorial_shots"

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0x00, 0x66, 0xCC)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_BG = RGBColor(0xF2, 0xF4, 0xF8)
CODE_BG = RGBColor(0x28, 0x2C, 0x34)
CODE_FG = RGBColor(0xE6, 0xE6, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_bg(slide, color=WHITE):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)
    return r


def add_kicker(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8), Inches(0.4))
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = "Segoe UI"


def add_title(slide, text, top=0.72, size=32):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = "Segoe UI"
    return tb


def add_rule(slide, top=1.55):
    ln = slide.shapes.add_connector(1, Inches(0.6), Inches(top), Inches(12.73), Inches(top))
    ln.line.color.rgb = ACCENT
    ln.line.width = Pt(2.25)


def add_bullets(slide, items, left=0.6, top=1.8, width=12.1, height=5.2, size=19, bold_first=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        run = p.add_run()
        run.text = ("•  " if level == 0 else "–  ") + text
        run.font.size = Pt(size if level == 0 else size - 2)
        run.font.color.rgb = NAVY if level == 0 else GRAY
        run.font.name = "Segoe UI"
        p.space_after = Pt(10 if level == 0 else 6)
    return tb


def add_code(slide, lines, left=0.6, top=1.9, width=12.1, height=None, size=15):
    height = height or (0.42 + 0.34 * len(lines))
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    box.adjustments[0] = 0.04
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = RGBColor(0x11, 0x13, 0x18)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.18)
    tf.margin_bottom = Inches(0.18)
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.name = "Consolas"
        run.font.color.rgb = CODE_FG
        p.space_after = Pt(2)
    return box


def add_image(slide, path, left, top, width=None, height=None, border=True):
    pic = slide.shapes.add_picture(path, Inches(left), Inches(top),
                                    width=Inches(width) if width else None,
                                    height=Inches(height) if height else None)
    if border:
        pic.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        pic.line.width = Pt(0.75)
    return pic


def add_footer(slide, text, page):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(7.08), Inches(9), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p.font.name = "Segoe UI"
    tb2 = slide.shapes.add_textbox(Inches(12.4), Inches(7.08), Inches(0.6), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = str(page)
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.RIGHT


def content_slide(kicker, title, page, footer="padb-tools"):
    s = add_slide()
    add_bg(s)
    add_kicker(s, kicker)
    add_title(s, title)
    add_rule(s)
    add_footer(s, footer, page)
    return s


# ---------------------------------------------------------------------------
# 1. Title
# ---------------------------------------------------------------------------
s = add_slide()
add_bg(s, NAVY)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.4))
p = tb.text_frame.paragraphs[0]
p.text = "padb-tools"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Segoe UI"

tb = s.shapes.add_textbox(Inches(0.9), Inches(3.75), Inches(11.5), Inches(0.9))
p = tb.text_frame.paragraphs[0]
p.text = "PADB Simple & PADB Interactive"
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(0x9F, 0xC5, 0xFF)
p.font.name = "Segoe UI"

tb = s.shapes.add_textbox(Inches(0.9), Inches(4.5), Inches(11.5), Inches(0.7))
p = tb.text_frame.paragraphs[0]
p.text = "Automating PADB-R.exe extraction and plotting for SG6311A characterization data"
p.font.size = Pt(17)
p.font.color.rgb = RGBColor(0xC4, 0xD2, 0xE8)
p.font.name = "Segoe UI"

tb = s.shapes.add_textbox(Inches(0.9), Inches(6.6), Inches(6), Inches(0.5))
p = tb.text_frame.paragraphs[0]
p.text = "SG6311A Test Engineering  |  C:\\apps\\padb\\tools"
p.font.size = Pt(13)
p.font.color.rgb = RGBColor(0x7E, 0x8F, 0xB3)

# ---------------------------------------------------------------------------
# 2. Agenda
# ---------------------------------------------------------------------------
s = content_slide("Overview", "Agenda", 2)
add_bullets(s, [
    "What is padb-tools?",
    "Three tiers: Legacy, Simple, Interactive",
    "PADB Simple — walkthrough",
    "PADB Interactive — walkthrough",
    "Setting up your environment to create your own plots",
    "Where to get help",
], size=22, top=2.0)

# ---------------------------------------------------------------------------
# 3. What is padb-tools
# ---------------------------------------------------------------------------
s = content_slide("Overview", "What is padb-tools?", 3)
add_bullets(s, [
    "Drives PADB-R.exe headlessly: runs a .pod file's analytics, collects the CSV/native output",
    "Turns results into interactive HTML — Plotly embedded inline, no server needed",
    "Works straight off a network share; results are self-contained files",
    "One job.json file configures a run — what pod, what date range, where to publish",
    ("PADB-R.exe is a real WinForms GUI app", 1),
    ("Needs an actual Windows desktop session — will not run over SSH or as a service", 1),
])

# ---------------------------------------------------------------------------
# 4. Three tiers
# ---------------------------------------------------------------------------
s = add_slide()
add_bg(s)
add_kicker(s, "Overview")
add_title(s, "Three Tiers — Pick What You Need")
add_rule(s)
add_footer(s, "padb-tools", 4)

cols = [
    ("Legacy", "mode omitted / \"legacy\"", [
        "Older per-analytic plot types",
        "accuracy_vs_freq, distribution,\nde_summary, etc.",
        "One job.json drives everything",
    ]),
    ("Simple", "\"mode\": \"simple\"", [
        "Literal PADB::Simple replacement",
        "PADB-R's own native PNG/PDF\nrenders, no custom plotting",
        "One job.json, same shape as Legacy",
    ]),
    ("Interactive", "\"mode\": \"interactive\"", [
        "Full modern interactive suite",
        "scatter, stat_summary, boxplot,\ndistribution, env_coverage, summary",
        "Two files: run job + plot job",
    ]),
]
col_w = 3.9
gap = 0.25
left0 = 0.6
for i, (name, key, bullets) in enumerate(cols):
    left = left0 + i * (col_w + gap)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.85),
                               Inches(col_w), Inches(4.9))
    card.adjustments[0] = 0.05
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = ACCENT
    card.line.width = Pt(1.25)
    card.shadow.inherit = False

    tb = s.shapes.add_textbox(Inches(left + 0.25), Inches(2.05), Inches(col_w - 0.5), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = name
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    tb = s.shapes.add_textbox(Inches(left + 0.25), Inches(2.6), Inches(col_w - 0.5), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = key
    p.font.size = Pt(13)
    p.font.name = "Consolas"
    p.font.color.rgb = ACCENT

    tb = s.shapes.add_textbox(Inches(left + 0.25), Inches(3.2), Inches(col_w - 0.5), Inches(3.3))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = "•  " + b
        r.font.size = Pt(14)
        r.font.color.rgb = GRAY
        p.space_after = Pt(14)

# ---------------------------------------------------------------------------
# 5. PADB Simple — what it is
# ---------------------------------------------------------------------------
s = content_slide("PADB Simple", "What It Is", 5)
add_bullets(s, [
    "A direct, static replacement for the old Perl PADB::Simple tool",
    "Zero custom plotting or statistics — no scatter, no tolerance intervals, no filters",
    "Drives PADB-R's own native rendering: OutputConfig_OutputGraph=1, GraphFormat=png,pdf",
    "Wraps the native PNG/PDF output in a bare gallery page, one card per image",
    "Metadata table dumped from the pod's own extraction settings",
    "Download links for the original .sao / .pod / .txt / .csv files",
    "Best for: a quick static report, or matching exactly what the legacy tool produced",
])

# ---------------------------------------------------------------------------
# 6. PADB Simple — job.json
# ---------------------------------------------------------------------------
s = content_slide("PADB Simple", "job.json", 6)
add_code(s, [
    "{",
    '  "description": "SG6311A MiniMoab Self-Test",',
    '  "pod": "SelfTestMiniMoab_000.pod",',
    '  "mode": "simple",',
    '  "padb_exe": "C:\\\\Program Files\\\\KEYSIGHT\\\\PADB-R.NET\\\\PADB-R.exe",',
    '  "results_dir": "SelfTestMiniMoab_000_simple_results",',
    '  "padb_timeout": 7200,',
    '  "run_analytics": true,',
    '  "publish": {',
    '    "destination": "\\\\\\\\srsnas01...\\\\SG6311A\\\\PADB-Simple\\\\MiniMoab\\\\SelfTestMiniMoab_000"',
    "  }",
    "}",
], size=14.5)

# ---------------------------------------------------------------------------
# 7. PADB Simple — running it
# ---------------------------------------------------------------------------
s = content_slide("PADB Simple", "Running It", 7)
add_bullets(s, [
    "Fastest path — generate the job.json straight from the pod:",
], top=1.85, height=0.5, size=19)
add_code(s, [
    "py padb_make_job.py MyPod.pod --module MyModule",
], top=2.35, height=0.55)
add_bullets(s, [
    "Then run it — real PADB-R.exe execution, needs a desktop session:",
], top=3.15, height=0.5, size=19)
add_code(s, [
    "py padb_run.py MyPod_job.json",
], top=3.65, height=0.55)
add_bullets(s, [
    "--module MyModule names the publish subfolder — always explicit, never guessed (more on slide 12)",
    "Output: results_dir\\index.html — one card per native PNG (linked to its PDF)",
    "HOW_TO_USE.txt is written alongside it, explaining what this tier does and doesn't do",
    "Re-plot without re-extracting: py padb_run.py MyPod_job.json --plots-only",
], top=4.35, height=2.3, size=18)

# ---------------------------------------------------------------------------
# 8. PADB Interactive — what it is
# ---------------------------------------------------------------------------
s = content_slide("PADB Interactive", "What It Is", 8)
add_bullets(s, [
    "The full modern interactive suite — six views from one Type=80 Scatter CSV",
    "scatter, stat_summary, boxplot, distribution, env_coverage, summary",
    "Auto-selects the right view set from the data itself:",
    ("Room-only data → scatter + boxplot", 1),
    ("Multi-temperature data → all six views", 1),
    "Live filters: condition, serial, port, frequency range, temperature",
    "Tolerance intervals (parametric TI, non-parametric TI), spec pass/fail",
    "Best for: any new pod — richer feature set, no regenerating to explore the data",
])

# ---------------------------------------------------------------------------
# 9. PADB Interactive — two-step workflow
# ---------------------------------------------------------------------------
s = add_slide()
add_bg(s)
add_kicker(s, "PADB Interactive")
add_title(s, "Two-Step Workflow")
add_rule(s)
add_footer(s, "padb-tools", 9)

steps = [
    ("1", "Extract", "padb_run.py", "Runs PADB-R.exe, writes a\nfresh CSV from the DB"),
    ("2", "Plot", "padb_v2.py", "Reads the CSV, builds all\nsix interactive HTML views"),
]
for i, (num, name, cmd, desc) in enumerate(steps):
    left = 0.9 + i * 6.0
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(2.0), Inches(0.9), Inches(0.9))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    tb = s.shapes.add_textbox(Inches(left + 1.1), Inches(2.0), Inches(4.5), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = name
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY

    add_code(s, [f"py {cmd} ..."], left=left, top=3.1, width=5.4, height=0.55, size=15)

    tb = s.shapes.add_textbox(Inches(left), Inches(3.85), Inches(5.4), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(15)
    p.font.color.rgb = GRAY

arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.35), Inches(2.2), Inches(0.5), Inches(0.5))
arrow.fill.solid()
arrow.fill.fore_color.rgb = RGBColor(0xCC, 0xD6, 0xE8)
arrow.line.fill.background()

tb = s.shapes.add_textbox(Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.9))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "CSV unchanged? Re-run step 2 alone — no need to touch PADB-R.exe again."
p.font.size = Pt(16)
p.font.italic = True
p.font.color.rgb = ACCENT

# ---------------------------------------------------------------------------
# 10. PADB Interactive — job.json (two files)
# ---------------------------------------------------------------------------
s = content_slide("PADB Interactive", "job.json — Two Files", 10)
add_bullets(s, ["*_run_job.json — same schema as any padb_run.py job:"], top=1.85, height=0.4, size=17)
add_code(s, [
    '{ "pod": "MyPod.pod", "mode": "interactive",',
    '  "results_dir": "my_analysis_run_results",',
    '  "run_analytics": true }',
], top=2.25, height=1.2, size=14)
add_bullets(s, ["*_v2_job.json — drives padb_v2.py, different schema:"], top=3.65, height=0.4, size=17)
add_code(s, [
    '{ "title_prefix": "SG6311A My Analysis",',
    '  "csv_path": "...\\\\my_analysis_run_results\\\\padb\\\\Scatter.csv",',
    '  "results_dir": "my_analysis_v2_results",',
    '  "spec_direction": "auto" }',
], top=4.05, height=1.5, size=14)
add_bullets(s, [
    "Fastest path — generate both at once from the pod:",
], top=5.85, height=0.4, size=17)
add_code(s, ["py padb_make_v2_job.py MyPod.pod --module MyModule"], top=6.25, height=0.55, size=15)

# ---------------------------------------------------------------------------
# 11. PADB Interactive — key features
# ---------------------------------------------------------------------------
s = content_slide("PADB Interactive", "Key Interactive Features", 11)
add_bullets(s, [
    "Filters — condition dropdowns, serial/DUT exclusion, port, frequency range, temperature",
    "Tolerance intervals — parametric TI and non-parametric (NP-TI), P/C selectors",
    "TLL display: Both / Upper only / Lower only",
    ("Shown live only when the CSV has no spec limit at all — a real limit always wins", 1),
    "Global Filter (GF) — set once in boxplot, propagates across every other view",
    "CSV export — download the currently filtered data from any view",
    "Everything recomputes client-side — no server, no regenerating to try a different filter",
])

# ---------------------------------------------------------------------------
# 12. Generators
# ---------------------------------------------------------------------------
s = content_slide("Both Tiers", "Generators — Skip Writing job.json by Hand", 12)
add_bullets(s, [
    "padb_make_job.py — Legacy / Simple / Interactive extract job, from a pod alone:",
], top=1.7, height=0.4, size=18)
add_code(s, ["py padb_make_job.py MyPod.pod --module MyModule --mode simple"],
         top=2.15, height=0.55, size=15)
add_bullets(s, [
    "padb_make_v2_job.py — full Interactive job set (run + one plot job per analytic):",
], top=2.85, height=0.4, size=18)
add_code(s, ["py padb_make_v2_job.py MyPod.pod --module MyModule"], top=3.3, height=0.55, size=15)
add_bullets(s, [
    "--module NAME sets the publish subfolder (e.g. MyModule → MiniMoab) — always explicit, never auto-guessed from the pod",
    "Both read the pod's own analytics directly — no manual guessing",
    "Auto-detects and fixes real pod issues: OutputFile collisions, disabled CSV output, one-sided specs with no configured limits",
    "Warns if a generated path nears Windows' 260-character limit",
], top=4.05, height=2.75, size=16.5)

# ---------------------------------------------------------------------------
# 13. Environment setup — prerequisites
# ---------------------------------------------------------------------------
s = content_slide("Environment Setup", "Prerequisites", 13)
add_bullets(s, [
    "Python 3 with the required packages:",
], top=1.85, height=0.4, size=19)
add_code(s, ["py -m pip install pandas numpy matplotlib scipy plotly"], top=2.25, height=0.55, size=15)
add_bullets(s, [
    "PADB-R.NET installed",
    ("Default path: C:\\Program Files\\KEYSIGHT\\PADB-R.NET\\PADB-R.exe", 1),
    ("Override with padb_exe in job.json if yours differs", 1),
    "A real Windows desktop session",
    ("PADB-R.exe is a WinForms app — needs an actual GUI session", 1),
    ("Will not run over SSH or as a headless service", 1),
    "A clone of the padb-tools repository",
], top=3.0, height=3.9, size=18)

# ---------------------------------------------------------------------------
# 14. Environment setup — per-user config
# ---------------------------------------------------------------------------
s = content_slide("Environment Setup", "Per-User Config (Optional, One-Time)", 14)
add_bullets(s, [
    "Every script derives defaults from Path.home() — never hardcoded to one username",
    "Want to customize? Create padb_config.json once:",
], top=1.85, height=0.9, size=19)
add_code(s, [
    "<Padb folder>\\padb_config.json",
    "{",
    '  "padb_exe": "C:\\\\Program Files\\\\KEYSIGHT\\\\PADB-R.NET\\\\PADB-R.exe",',
    '  "padb_output_dir": "...\\\\Padb\\\\R-Plots",',
    '  "padb_logs_dir": "...\\\\Padb\\\\Logs",',
    '  "data_dir": "...\\\\Padb\\\\Data",',
    '  "publish_root": "\\\\\\\\srsnas01...\\\\SG6311A\\\\PADB-Simple"',
    "}",
], top=2.75, height=2.85, size=13.5)
add_bullets(s, [
    "Entirely optional — skip it and every key falls back to a sensible default",
], top=5.85, height=0.5, size=18)

# ---------------------------------------------------------------------------
# 15. Environment setup — first run checklist
# ---------------------------------------------------------------------------
s = content_slide("Environment Setup", "Your First Plot — Checklist", 15)
add_bullets(s, [
    "Get your .pod file",
    "Generate a job.json (or two, for Interactive):",
    ("padb_make_job.py MyPod.pod --module MyModule --mode simple", 1),
    ("padb_make_v2_job.py MyPod.pod --module MyModule", 1),
    "Run the extraction: py padb_run.py <job.json>",
    "Interactive only — build the plots: py padb_v2.py <v2_job.json>",
    "Open results_dir\\index.html",
    "Before publishing, sanity-check against QA_Checklist.md",
], size=20)

# ---------------------------------------------------------------------------
# 16. Where to get help
# ---------------------------------------------------------------------------
s = content_slide("Wrap-Up", "Where to Get Help", 16)
add_bullets(s, [
    "New to this? Read in order:",
    ("GETTING_STARTED.md → Quick_Start.md → PADB_Tools_Guide.md", 1),
    "One-page steps, nothing else:",
    ("Simple_Mode_Cheatsheet.md  /  Interactive_Mode_Cheatsheet.md", 1),
    "Authoring or fixing a pod:",
    ("PADB_Analytic_Requirements.md", 1),
    "Before you publish new results:",
    ("QA_Checklist.md", 1),
    "Using Claude Code? Type /padb-tools for an on-demand summary of the",
    ("architecture, gotchas, and new-pod checklist — no need to open all six docs", 1),
], size=18)

# ---------------------------------------------------------------------------
# 17. Environment setup — PATH vs PYTHONPATH
# ---------------------------------------------------------------------------
s = content_slide("Environment Setup", "Making Python Find Your Script Folder", 17)
add_bullets(s, [
    "Two different goals need two different environment variables:",
], top=1.85, height=0.4, size=19)
add_bullets(s, [
    "Run scripts by name from any directory (padb_run.py job.json):",
], top=2.35, height=0.4, size=18)
add_code(s, ['setx PATH "$env:PATH;C:\\apps\\padb\\tools"'], top=2.75, height=0.55, size=15)
add_bullets(s, [
    ("Open a NEW terminal afterward — only helps for bare invocation", 1),
    ("(padb_run.py ...); typing \"py padb_run.py ...\" still resolves the", 1),
    ("path relative to your current directory, PATH doesn't help there", 1),
], top=3.4, height=1.0, size=14.5)
add_bullets(s, [
    "Let other scripts import padb_run / padb_config from anywhere:",
], top=4.5, height=0.4, size=18)
add_code(s, ['setx PYTHONPATH "C:\\apps\\padb\\tools"'], top=4.9, height=0.55, size=15)
add_bullets(s, [
    ("Already have entries? Append instead: setx PYTHONPATH", 1),
    ('("$env:PYTHONPATH;C:\\apps\\padb\\tools")', 1),
], top=5.55, height=0.7, size=14.5)

# ---------------------------------------------------------------------------
# 18. Tutorial 1 intro — PADB Simple
# ---------------------------------------------------------------------------
s = content_slide("Tutorial 1 — PADB Simple", "The Pod: MaxPowerTutorial1.pod", 18)
add_bullets(s, [
    "A real MaxPower-family pod, 6 analytics:",
    ("4 x Type=80 Scatter — Unleveled Log, Leveled Log, Leveled Linear, Unleveled Linear", 1),
    ("1 x Type=90 SummaryPlot — Summary_Leveled_Linear", 1),
    ("1 x Type=60 Environmental — Environmental_Leveled_Linear", 1),
    "No spec limits configured (Limits_YLimit=None on every analytic) — Simple mode doesn't",
    ("care; it posts whatever PADB-R renders natively, spec or no spec", 1),
    "Goal: turn this pod into a browsable gallery with zero custom plotting code",
], size=18)

# ---------------------------------------------------------------------------
# 19. Tutorial 1 — Step 1: drop & generate (GUI)
# ---------------------------------------------------------------------------
s = content_slide("Tutorial 1 — PADB Simple", "Step 1 — Drop the Pod, Generate the Job", 19)
add_image(s, f"{IMG_DIR}/t1_step1.png", left=0.6, top=1.7, width=7.6)
tb = s.shapes.add_textbox(Inches(8.5), Inches(2.0), Inches(4.3), Inches(4.5))
tf = tb.text_frame
tf.word_wrap = True
for i, txt in enumerate([
    "py webapp\\padb_web.py",
    "Drag MaxPowerTutorial1.pod onto the page",
    "Analytics parse instantly — no extraction needed to preview them",
    "Mode: simple, Module: Tutorial",
    "Click Generate Job — writes MaxPowerTutorial1_job.json next to the pod",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(15 if i else 14)
    r.font.bold = (i == 0)
    r.font.color.rgb = NAVY if i == 0 else GRAY
    p.space_after = Pt(10)

# ---------------------------------------------------------------------------
# 20. Tutorial 1 — Step 2: run it (GUI)
# ---------------------------------------------------------------------------
s = content_slide("Tutorial 1 — PADB Simple", "Step 2 — Check the Box, Run It", 20)
add_image(s, f"{IMG_DIR}/t1_step2.png", left=0.6, top=1.7, width=8.6)
tb = s.shapes.add_textbox(Inches(9.5), Inches(2.0), Inches(3.3), Inches(4.5))
tf = tb.text_frame
tf.word_wrap = True
for i, txt in enumerate([
    "Check the job's row, click Run Selected",
    "Real PADB-R.exe runs — the queue serializes it, so it's safe even if others are running jobs too",
    "Status panel streams the live log",
    "\"done\" + Open results link when it finishes — 107.1s, this run",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(15 if i == 0 else 14)
    r.font.bold = (i == 0)
    r.font.color.rgb = NAVY if i == 0 else GRAY
    p.space_after = Pt(10)

# ---------------------------------------------------------------------------
# 21. Tutorial 1 — the result
# ---------------------------------------------------------------------------
s = content_slide("Tutorial 1 — PADB Simple", "The Result", 21)
add_image(s, f"{IMG_DIR}/simple_index.png", left=0.6, top=1.75, width=7.2)
tb = s.shapes.add_textbox(Inches(9.5), Inches(2.0), Inches(3.3), Inches(4.5))
tf = tb.text_frame
tf.word_wrap = True
for i, txt in enumerate([
    "results_dir\\index.html",
    "One card per native PNG, largest first",
    "Metadata table (collapsed) shows the pod's own extraction settings",
    "Download links for .sao / .pod / .txt / .csv",
    "Zero custom code involved — this is exactly what PADB-R rendered",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run()
    r.text = ("•  " if i else "") + txt if i else txt
    r.font.size = Pt(15 if i == 0 else 14)
    r.font.bold = (i == 0)
    r.font.color.rgb = NAVY if i == 0 else GRAY
    p.space_after = Pt(10)

# ---------------------------------------------------------------------------
# 22. Tutorial 2 intro — PADB Interactive
# ---------------------------------------------------------------------------
s = content_slide("Tutorial 2 — PADB Interactive", "The Pod: MaxPowerTutorial2.pod", 22)
add_bullets(s, [
    "A focused MaxPower-family pod, 2 analytics — both Type=80 Scatter:",
    ("Leveled Linear, Unleveled Linear", 1),
    "No spec limits configured — same as Tutorial 1's pod",
    "Fewer analytics than Tutorial 1 on purpose: Interactive generates a full",
    ("6-view suite per analytic, so 2 analytics already means 12 HTML files", 1),
    "Goal: the full interactive suite — filters, tolerance intervals, the live",
    ("TLL display selector — not just a static gallery", 1),
], size=18)

# ---------------------------------------------------------------------------
# 23. Tutorial 2 — Step 1: generate both job files
# ---------------------------------------------------------------------------
s = content_slide("Tutorial 2 — PADB Interactive", "Step 1 — Drop the Pod, Generate the Job", 23)
add_image(s, f"{IMG_DIR}/t2_step1.png", left=0.6, top=1.7, width=7.6)
tb = s.shapes.add_textbox(Inches(8.5), Inches(2.0), Inches(4.3), Inches(4.5))
tf = tb.text_frame
tf.word_wrap = True
for i, txt in enumerate([
    "Same drop, different mode",
    "Mode: interactive (V2), Module: Tutorial",
    "One click, three files: 1 run job (extraction) + 1 plot job per analytic",
    "csv_path in each plot job is PREDICTED — padb_v2.py verifies it against the real CSV once the run job has actually extracted",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(15 if i == 0 else 14)
    r.font.bold = (i == 0)
    r.font.color.rgb = NAVY if i == 0 else GRAY
    p.space_after = Pt(10)

# ---------------------------------------------------------------------------
# 24. Tutorial 2 — Step 2: select just the run job
# ---------------------------------------------------------------------------
s = content_slide("Tutorial 2 — PADB Interactive", "Step 2 — Select Just the Run Job", 24)
add_bullets(s, [
    "Three rows appear — the run job and both plot jobs. Check only the run job:",
], top=1.75, height=0.5, size=17)
add_image(s, f"{IMG_DIR}/t2_step2_top.png", left=0.6, top=2.35, width=12.1)
add_bullets(s, [
    "You never need to check the plot-job rows yourself for a fresh dataset — the run job finds and runs them automatically once extraction succeeds",
], top=6.05, height=0.7, size=15.5)

# ---------------------------------------------------------------------------
# 25. Tutorial 2 — Step 3: one click, extraction + all plots
# ---------------------------------------------------------------------------
s = content_slide("Tutorial 2 — PADB Interactive", "Step 3 — One Click: Extract AND Build Plots", 25)
add_bullets(s, [
    "Click Run Selected once. Real PADB-R.exe extraction, then both plot jobs build automatically — no second step:",
], top=1.65, height=0.5, size=15.5)
add_image(s, f"{IMG_DIR}/t2_step2_bottom.png", left=2.15, top=2.2, width=9.0)
add_bullets(s, [
    "3 temperatures (Room, 20°C, 30°C) — padb_v2.py auto-detects multi-temp data and builds the full six-view suite, no \"views\" key needed",
], top=6.35, height=0.7, size=13.5)

# ---------------------------------------------------------------------------
# 26. Tutorial 2 — the result: combined gallery
# ---------------------------------------------------------------------------
s = content_slide("Tutorial 2 — PADB Interactive", "The Result — Combined Gallery", 26)
add_image(s, f"{IMG_DIR}/interactive_index.png", left=0.9, top=1.85, width=6.0)
tb = s.shapes.add_textbox(Inches(7.3), Inches(2.1), Inches(5.4), Inches(4.5))
tf = tb.text_frame
tf.word_wrap = True
for i, txt in enumerate([
    "MaxPowerTutorial2_v2_results\\index.html",
    "One index, 12 links: 2 analytics x 6 views each",
    "(boxplot, distribution, env coverage, scatter, stat summary, summary)",
    "Re-run just padb_v2.py to rebuild this instantly if you tweak a job.json —",
    "no need to touch PADB-R.exe again unless the CSV itself changed",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(15 if i == 0 else 14)
    r.font.bold = (i == 0)
    r.font.color.rgb = NAVY if i == 0 else GRAY
    p.space_after = Pt(10)

# ---------------------------------------------------------------------------
# 27. Tutorial 2 — the result: summary view + TLL selector
# ---------------------------------------------------------------------------
s = content_slide("Tutorial 2 — PADB Interactive", "The Result — Try the TLL Selector", 27)
add_bullets(s, [
    "This pod has no spec limits, so the live \"TLL display\" selector appears (circled) —",
    "MaxPower is a guaranteed-minimum-power measurement, so try switching it to \"Lower only\"",
], top=1.75, height=0.8, size=15)
_img_left, _img_top, _img_w = 2.07, 2.65, 9.2
add_image(s, f"{IMG_DIR}/interactive_summary_crop.png", left=_img_left, top=_img_top, width=_img_w)
_hl = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(_img_left + 0.418 * _img_w - 0.05), Inches(_img_top + 0.154 * _img_w * 0.4643 - 0.05),
                          Inches(0.671 * _img_w - 0.418 * _img_w + 0.1), Inches(0.188 * _img_w * 0.4643 - 0.154 * _img_w * 0.4643 + 0.1))
_hl.fill.background()
_hl.line.color.rgb = RGBColor(0xCC, 0x00, 0x00)
_hl.line.width = Pt(2.25)
_hl.shadow.inherit = False

# ---------------------------------------------------------------------------
# 27b. Interactive deep dive — the six views at a glance
# ---------------------------------------------------------------------------
s = content_slide("PADB Interactive — Deep Dive", "The Six Views at a Glance", 28)
grid = [
    ("t2_view_scatter.png", "Scatter — raw per-measurement, all temps"),
    ("t2_view_boxplot.png", "Boxplot — box stats per condition/temperature"),
    ("t2_view_stat_summary.png", "Stat Summary — room-temp per-DUT population stats"),
    ("t2_view_summary.png", "Summary — all-temp NP-TI summary"),
    ("t2_view_env_coverage.png", "Env Coverage — carrier-power stability vs. temperature"),
    ("t2_view_distribution.png", "Distribution — multi-temp Delta-Env KDE (Harmonics pod)"),
]
gw, gx0, gy0, gpx, gpy = 3.6, 0.6, 1.65, 0.2, 0.4
img_h = gw * (950 / 1600)
for i, (fname, cap) in enumerate(grid):
    r, c = divmod(i, 3)
    left = gx0 + c * (gw + gpx)
    top = gy0 + r * (img_h + gpy)
    add_image(s, f"{IMG_DIR}/{fname}", left=left, top=top, width=gw)
    tb = s.shapes.add_textbox(Inches(left), Inches(top + img_h + 0.05), Inches(gw), Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = cap
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY

# ---------------------------------------------------------------------------
# 27c. Interactive deep dive — segment tab-through
# ---------------------------------------------------------------------------
s = content_slide("PADB Interactive — Deep Dive", "Segment Tab-Through — Jump Between Spec Bands", 29)
add_image(s, f"{IMG_DIR}/deepdive_segment_tabs.png", left=0.55, top=1.7, width=8.0)
tb = s.shapes.add_textbox(Inches(8.85), Inches(2.0), Inches(4.0), Inches(4.8))
tf = tb.text_frame
tf.word_wrap = True
for i, txt in enumerate([
    "ClockSpurs_PADBToolTest.pod",
    "A real 5-level spec staircase per SpurType",
    "\"Segment by\": Spec / Limit / Uncertainty",
    "Prev / Next jump the freq window to each contiguous spec band",
    "Isolate one SpurType first (dropdown, top-left) — segments are per-selection",
    "Respects the Global Filter and serial/port filters: exclude a DUT and its own spec/limit/uncertainty contribution drops too",
    "All six views: scatter, boxplot, distribution already had this; stat_summary, summary, env_coverage, and the real Delta-Env distribution view all got it this session",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(15 if i == 0 else 13.5)
    r.font.bold = (i == 0)
    r.font.color.rgb = NAVY if i == 0 else GRAY
    p.space_after = Pt(9)

# ---------------------------------------------------------------------------
# 27d. Interactive deep dive — Group By (stat_summary)
# ---------------------------------------------------------------------------
s = content_slide("PADB Interactive — Deep Dive", "Group By — Collapsing Fragmented Conditions", 30)
add_image(s, f"{IMG_DIR}/deepdive_groupby_stat_summary.png", left=0.55, top=1.7, width=8.4)
tb = s.shapes.add_textbox(Inches(9.2), Inches(2.0), Inches(3.65), Inches(4.8))
tf = tb.text_frame
tf.word_wrap = True
for i, txt in enumerate([
    "Stat Summary — Group by: SpurType",
    "This pod's per-unit Group text fragments into 151 near-duplicate conditions",
    "Group by collapses them to the 5 real spur types — one legend entry, one spec line each",
    "Mean/min/max/quantiles pool exactly per DUT — a DUT's data falls under exactly one condition per dimension",
    "NP-TI / spec take the worst case across the collapsed conditions",
    "Same dropdown added to Summary and Env Coverage this session",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(15 if i == 0 else 13.5)
    r.font.bold = (i == 0)
    r.font.color.rgb = NAVY if i == 0 else GRAY
    p.space_after = Pt(9)

# ---------------------------------------------------------------------------
# 27e. Interactive deep dive — Group By + segment tabs (env_coverage)
# ---------------------------------------------------------------------------
s = content_slide("PADB Interactive — Deep Dive", "Environmental Coverage — Group By + Segment Tabs", 31)
add_image(s, f"{IMG_DIR}/deepdive_groupby_env_coverage.png", left=0.55, top=1.7, width=8.4)
tb = s.shapes.add_textbox(Inches(9.2), Inches(2.0), Inches(3.65), Inches(4.8))
tf = tb.text_frame
tf.word_wrap = True
for i, txt in enumerate([
    "Env Coverage — Group by: SpurType, Segment by: Spec",
    "Before this session: one flat spec value per condition, no per-frequency tracking at all",
    "Now: per-frequency Upper/Lower Limit, Spec, and Uncertainty tracked and segment-tabbable, same as every other view",
    "UDE/LDE/TTU/TTL recompute exactly on every Group By change — already derived from raw per-DUT data on each call",
    "GF toggle (top-right) still applies per-DUT exclusions the same way",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(15 if i == 0 else 13.5)
    r.font.bold = (i == 0)
    r.font.color.rgb = NAVY if i == 0 else GRAY
    p.space_after = Pt(9)

# ---------------------------------------------------------------------------
# 28. Tutorial 3 intro — /padb-tools
# ---------------------------------------------------------------------------
s = content_slide("Tutorial 3 — /padb-tools", "What It Is, When to Use It", 32)
add_bullets(s, [
    "A Claude Code command scoped to this repo (.claude/commands/padb-tools.md)",
    "Surfaces a condensed index of the tool's architecture, gotchas, and a",
    ("13-point new-pod checklist — without opening all six markdown docs", 1),
    "Only works when Claude Code's working directory is inside the repo",
    ("(e.g. C:\\apps\\padb\\tools) — that's what makes it \"repo-scoped\"", 1),
    "Use it when you're: onboarding a new pod, debugging a plot feature,",
    ("reviewing generated HTML, or just orienting yourself in a new session", 1),
], size=18)

# ---------------------------------------------------------------------------
# 29. Tutorial 3 — what it actually returns
# ---------------------------------------------------------------------------
s = content_slide("Tutorial 3 — /padb-tools", "What It Actually Tells You", 33)
add_bullets(s, [
    "A few of the 13 new-pod checklist items it surfaces immediately:",
    ("TestRun_RunStatus — exit 0 but no CSV? Check the pod's default 'P'-only filter", 1),
    ("Environment_TestStep must be {All} or distribution/env_coverage/summary go empty", 1),
    ("OutputConfig_OutputCSV=False on a Scatter analytic — renders PNG, writes no CSV", 1),
    ("Generated job.json path length — warns before Windows' 260-char limit bites", 1),
    "Plus: the architecture one-liner, common JS gotchas, and pointers to the",
    ("full doc for anything it only summarizes", 1),
], size=17)

# ---------------------------------------------------------------------------
# 30. Tutorial 3 — worked example
# ---------------------------------------------------------------------------
s = content_slide("Tutorial 3 — /padb-tools", "Worked Example: MaxPowerTutorial3.pod", 34)
add_bullets(s, [
    "You've just been handed MaxPowerTutorial3.pod — 2 analytics, never run before",
    "Open Claude Code with your working directory inside the padb-tools repo",
    "Type: /padb-tools",
    "Ask directly: \"Is MaxPowerTutorial3.pod ready to run in Interactive mode?\"",
    "Claude reads the checklist above, then actually inspects the pod file —",
    ("Limits_YLimit, Environment_TestStep, OutputConfig_OutputCSV, Group strings", 1),
    "You get a direct answer grounded in this pod's real settings, not a guess",
], size=18)

# ---------------------------------------------------------------------------
# 30b. Site conversion — the problem
# ---------------------------------------------------------------------------
s = content_slide("Multi-Site Testing", "Same Test, Different Database", 35)
add_bullets(s, [
    "Malaysia (AMC2) production ramp-up: the same test now also pulls from",
    ("AMC2's own PADB Oracle database, not just Santa Rosa's", 1),
    "Comparing a real Santa Rosa pod against its AMC2 counterpart: the only",
    ("genuine differences are Device_Server and Device_Database in [Extract]", 1),
    "Every AnalyticName and OutputConfig_OutputFile is otherwise identical —",
    ("running both site variants writes identically-named CSVs", 1),
    "Risk: point both at a shared results/publish location and one silently",
    ("overwrites the other — same class of problem unique_output_filenames", 1),
    ("solves within one pod, but across two site-variant pods instead", 1),
], size=17)

# ---------------------------------------------------------------------------
# 30c. Site conversion — converting a pod
# ---------------------------------------------------------------------------
s = content_slide("Multi-Site Testing", "padb_convert_site.py — Converting a Pod", 36)
add_bullets(s, ["Site registry lives in padb_sites.json — add a new site there, no code changes:"],
            top=1.75, height=0.4, size=16)
add_code(s, [
    "py padb_convert_site.py --pod MyPod.pod --to AMC2",
], top=2.15, height=0.55, size=14.5)
add_code(s, [
    "Wrote MyPod-AMC2.pod  (SantaRosa -> AMC2)",
    "  Device_Server -> PADB ORACLE AMC2",
    "  Device_Database -> GALLEON_1",
    "  AnalyticName 'Leveled Linear' -> 'Leveled Linear AMC2'",
    "  OutputConfig_OutputFile '..._Linear' -> '..._Linear_AMC2'",
    "  WARNING: SaoFile now points at 'MyPod-AMC2.sao', which does not",
    "  exist yet -- supply a real .sao extracted at AMC2 before running",
], top=2.85, height=2.5, size=13.5)
add_bullets(s, [
    "Source pod is never touched. Round-trip verified byte-identical",
    ("(SantaRosa -> AMC2 -> SantaRosa reproduces the original exactly)", 1),
], top=5.55, height=0.7, size=15)

# ---------------------------------------------------------------------------
# 30d. Site conversion — converting a job
# ---------------------------------------------------------------------------
s = content_slide("Multi-Site Testing", "padb_convert_site.py — Converting a Job", 37)
add_code(s, [
    "py padb_convert_site.py --job my_run_job.json --to AMC2",
], top=1.85, height=0.55, size=14.5)
add_bullets(s, ["Repoints \"pod\" and swaps the pod stem everywhere it appears:"],
            top=2.6, height=0.4, size=17)
add_code(s, [
    "Converted pod not found yet -- creating it first: ...",
    "Wrote MyPod-AMC2_run_job.json  (SantaRosa -> AMC2,",
    "  pod=MyPod-AMC2.pod)",
    "  This is a V2 run job -- also regenerate its plot",
    "  job(s) against the new pod:",
    "    py padb_make_v2_job.py MyPod-AMC2.pod --module <YourModule>",
], top=3.0, height=2.3, size=13.5)
add_bullets(s, [
    "Updates results_dir, publish/publish_to, and description too — not just \"pod\"",
    "--force overwrites an existing output; --list-sites prints the registry",
], top=5.55, height=0.7, size=15)

# ---------------------------------------------------------------------------
# 31. Thank you
# ---------------------------------------------------------------------------
s = add_slide()
add_bg(s, NAVY)
tb = s.shapes.add_textbox(Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.2))
p = tb.text_frame.paragraphs[0]
p.text = "Questions?"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE

tb = s.shapes.add_textbox(Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.7))
p = tb.text_frame.paragraphs[0]
p.text = "C:\\apps\\padb\\tools   |   /padb-tools"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0x9F, 0xC5, 0xFF)
p.font.name = "Consolas"

# ---------------------------------------------------------------------------
# 32. Appendix — divider
# ---------------------------------------------------------------------------
s = add_slide()
add_bg(s, LIGHT_BG)
tb = s.shapes.add_textbox(Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.0))
p = tb.text_frame.paragraphs[0]
p.text = "Appendix"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = NAVY
tb = s.shapes.add_textbox(Inches(0.9), Inches(3.9), Inches(11.5), Inches(0.7))
p = tb.text_frame.paragraphs[0]
p.text = "Setting Up From Zero — Assumes No Prior Access"
p.font.size = Pt(22)
p.font.color.rgb = ACCENT

# ---------------------------------------------------------------------------
# 33. Appendix — getting the code
# ---------------------------------------------------------------------------
s = content_slide("Appendix", "Getting the Code", 39)
add_bullets(s, [
    "This repo isn't on a company-wide share by default — copy it from:",
], top=1.85, height=0.4, size=19)
add_code(s, [
    "\\\\srsnas01.srs.is.keysight.com\\prod\\MIDRF3\\SG6311A\\Padb-tools\\tools",
], top=2.3, height=0.55, size=13.5)
add_bullets(s, [
    "Copy that whole \"tools\" folder to a local path, e.g. C:\\apps\\padb\\tools",
], top=3.05, height=0.4, size=18)
add_bullets(s, [
    "padb_batch.py is bundled directly inside the tools folder (added 2026-08-05) —",
    ("no separate dependency to track down; the whole repo is self-contained", 1),
], top=3.75, height=1.0, size=17)
add_bullets(s, [
    "Then follow the Prerequisites and Per-User Config slides earlier in this deck",
], top=4.9, height=0.5, size=16.5)

# ---------------------------------------------------------------------------
# 34. Appendix — first commands to verify
# ---------------------------------------------------------------------------
s = content_slide("Appendix", "First Commands to Verify Your Setup", 40)
add_bullets(s, ["From inside the tools folder, confirm everything imports cleanly:"],
            top=1.85, height=0.4, size=18)
add_code(s, ['py -c "import padb_run, padb_v2, padb_config"'], top=2.3, height=0.55, size=14.5)
add_bullets(s, [
    "No output and no traceback = success",
], top=3.0, height=0.4, size=17)
add_bullets(s, ["Then generate your first job.json from any pod you have access to:"],
            top=3.6, height=0.4, size=18)
add_code(s, ["py padb_make_job.py YourPod.pod --module YourModule"], top=4.05, height=0.55, size=14.5)
add_bullets(s, [
    "If PADB-R.exe launches and a real window briefly appears, you're fully set up",
    "Stuck? GETTING_STARTED.md → Quick_Start.md, or type /padb-tools once you're",
    ("working inside the repo", 1),
], top=4.75, height=1.6, size=17)

prs.save("PADB_Simple_and_Interactive.pptx")
print(f"Saved PADB_Simple_and_Interactive.pptx -- {len(prs.slides._sldIdLst)} slides")
